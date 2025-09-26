import os
import gradio as gr
from typing import Optional, List, Tuple, Union
import time
import base64
import json
from PIL import Image as PILImage
import io
from pathlib import Path

# File readers
from pypdf import PdfReader
import docx
from pptx import Presentation

# Import your existing helper functions
from brain_of_the_doctor import encode_image, analyze_image_with_query
from voice_of_the_doctor import text_to_speech_with_gtts
from voice_of_the_patinet import transcribe_with_groq
from langchain_doctor import DoctorChain

# System prompt for doctor chain
SYSTEM_PROMPT = """You have to act as a professional doctor, i know you are not but this is for learning purpose. 
What's in this image?. Do you find anything wrong with it medically? 
If you make a differential, suggest some remedies for them. Donot add any numbers or special characters in 
your response. Your response should be in one long paragraph. Also always answer as if you are answering to a real person.
Donot say 'In the image I see' but say 'With what I see, I think you have ....' only for the first prompt
Dont respond as an AI model in markdown, your answer should mimic that of an actual doctor not an AI bot, 
Keep your answer concise (max 2 sentences). No preamble, start your answer right away please"""

# -------------------------
# File reading helper
# -------------------------
def read_file_content(file_path: str) -> str:
    """Extract text from PDF, DOCX, PPTX, TXT, etc."""
    ext = Path(file_path).suffix.lower()
    text = ""

    try:
        if ext == ".pdf":
            reader = PdfReader(file_path)
            for page in reader.pages:
                if page.extract_text():
                    text += page.extract_text() + "\n"

        elif ext == ".docx":
            doc = docx.Document(file_path)
            text = "\n".join([para.text for para in doc.paragraphs])

        elif ext == ".pptx":
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"

        elif ext in [".txt", ".md", ".csv"]:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()

        else:
            return "Unsupported file format."

    except Exception as e:
        return f"Error reading file: {e}"

    return text.strip() if text else "No readable text found in the file."


class AIDoctor:
    def __init__(self):
        self.doctor_chain = DoctorChain()
        self.conversation_history = []

    def image_to_base64(self, image_path: str) -> str:
        """Convert image to base64 for HTML display"""
        try:
            with open(image_path, "rb") as img_file:
                img_data = img_file.read()
                base64_str = base64.b64encode(img_data).decode()
                return f"data:image/jpeg;base64,{base64_str}"
        except Exception as e:
            print(f"Error converting image to base64: {e}")
            return ""

    def process_message(
        self,
        message: str,
        chat_history_html: str,
        audio_filepath: Optional[str] = None,
        image_filepath: Optional[str] = None,
        file_filepath: Optional[str] = None
    ) -> Tuple[str, str, Optional[str]]:
        """Process user message and return updated HTML chat"""
        try:
            user_input = ""
            file_content = ""
            
            # Handle audio input
            if audio_filepath:
                try:
                    user_input = transcribe_with_groq(
                        GROQ_API_KEY=os.environ.get("GROQ_API_KEY"),
                        audio_filepath=audio_filepath,
                        stt_model="whisper-large-v3"
                    )
                    print(f"Transcribed audio: {user_input}")
                except Exception as e:
                    print(f"Audio transcription error: {e}")
                    user_input = message
            else:
                user_input = message

            if not user_input.strip() and not image_filepath:
                return "", chat_history_html, None

            # Handle file input
            if file_filepath:
                file_content = read_file_content(file_filepath)
                #user_input += f"\n\n[Patient uploaded a file, extracted content:]\n{file_text}"

            if not user_input.strip() and not image_filepath and not file_filepath:
                return "", chat_history_html, None

            if image_filepath:
                try:
                    encoded_image = encode_image(image_filepath)
                    
                    query_with_file = SYSTEM_PROMPT + "\nPatient's Query: " + user_input
                    if file_content:
                        query_with_file += f"\n\n[Patient uploaded file content:]\n{file_content}"

                    doctor_response = analyze_image_with_query(
                        #query=SYSTEM_PROMPT + "\nPatient's Query: " + user_input,
                        query=query_with_file,
                        encoded_image=encoded_image,
                        model="meta-llama/llama-4-scout-17b-16e-instruct"
                    )
                    print(f"Image analysis response: {doctor_response}")

                    # ✅ Manually save this turn into memory
                    self.doctor_chain.save_to_memory(user_input, doctor_response)

                except Exception as e:
                    print(f"Image analysis error: {e}")
                    doctor_response = "I couldn't analyze the image properly, please try again."
            else:
                try:
                    final_query = f"Patient said: {user_input}"
                    
                    if file_content:
                        final_query += f"\n\n[Patient uploaded file content:]\n{file_content}"
                    
                    doctor_response = self.doctor_chain.get_response(query=final_query)
                    print(f"Doctor response: {doctor_response}")

                    # ✅ Already auto-saves inside get_response,
                    # but you can call again explicitly if you want:
                    # self.doctor_chain.save_to_memory(final_query, doctor_response)

                except Exception as e:
                    print(f"Doctor chain error: {e}")
                    doctor_response = "I'm having trouble processing your request at the moment. Please try again."


            # Generate audio response
            audio_output_path = None
            try:
                output_path = f"response_{int(time.time())}.mp3"
                text_to_speech_with_gtts(
                    input_text=doctor_response,
                    output_filepath=output_path
                )
                audio_output_path = output_path
                print(f"Generated audio: {audio_output_path}")
            except Exception as e:
                print(f"TTS Error: {e}")

            # Build HTML for the new messages
            new_chat_html = self.build_chat_html(
                user_input, doctor_response, image_filepath, file_filepath, chat_history_html
            )

            return "", new_chat_html, audio_output_path

        except Exception as e:
            error_msg = f"Sorry, I encountered an error: {str(e)}"
            print(f"Process message error: {e}")
            
            new_chat_html = self.build_chat_html(
                user_input if 'user_input' in locals() else message, 
                error_msg, 
                image_filepath, 
                file_filepath if 'file_filepath' in locals() else None,
                chat_history_html
            )
            
            return "", new_chat_html, None

    def build_chat_html(self, user_message: str, ai_response: str, image_path: str, file_path: str, current_html: str) -> str:
        """Build HTML for chat messages with ChatGPT-like styling"""
        
        # ✅ NEW: File display logic
        file_html = ""
        if file_path:
            file_name = Path(file_path).name
            file_ext = Path(file_path).suffix.lower()
            
            # Choose appropriate icon based on file type
            if file_ext == ".pdf":
                file_icon = "📄"
            elif file_ext in [".docx", ".doc"]:
                file_icon = "📝"
            elif file_ext in [".pptx", ".ppt"]:
                file_icon = "📊"
            elif file_ext in [".txt", ".md"]:
                file_icon = "📋"
            elif file_ext == ".csv":
                file_icon = "📊"
            else:
                file_icon = "📎"
            
            file_html = f"""
            <div class="file-attachment">
                <span class="file-icon">{file_icon}</span>
                <span class="file-name">{file_name}</span>
            </div>
            """
        
        # User message HTML
        user_html = f"""
        <div class="message user-message">
            <div class="message-content">
                {"<img src='" + self.image_to_base64(image_path) + "' class='message-image'>" if image_path else ""}
                {file_html}
                <div class="message-text">{user_message}</div>
            </div>
        </div>
        """
        
        # AI response HTML
        ai_html = f"""
        <div class="message ai-message">
            <div class="avatar">🩺</div>
            <div class="message-content">
                <div class="message-text">{ai_response}</div>
            </div>
        </div>
        """
        
        return current_html + user_html + ai_html

    def clear_conversation(self):
        """Clear the conversation"""
        self.doctor_chain = DoctorChain()
        return "", ""


def create_interface():
    """Create ChatGPT-like interface using HTML"""
    
    ai_doctor = AIDoctor()
    
    # CSS for consistent dark theme matching the image
    custom_css = """
    <style>
        /* Global dark theme override */
        .gradio-container {
            background-color: #1a1a1a !important;
        }
        
        body {
            background-color: #1a1a1a !important;
        }
        
        /* Dark theme for all Gradio components */
        .gradio-container .block {
            background-color: #1a1a1a !important;
            border-color: #404040 !important;
        }
        
        .gradio-container .form {
            background-color: #1a1a1a !important;
        }
        
        .gradio-container label {
            color: #ffffff !important;
        }
        
        .gradio-container input, .gradio-container textarea {
            background-color: #404040 !important;
            border-color: #555 !important;
            color: #ffffff !important;
        }
        
        .gradio-container input::placeholder, .gradio-container textarea::placeholder {
            color: #aaa !important;
        }
        
        .gradio-container .markdown {
            color: #ffffff !important;
        }
        
        .gradio-container button {
            background-color: #404040 !important;
            border-color: #555 !important;
            color: #ffffff !important;
        }
        
        .gradio-container button:hover {
            background-color: #555 !important;
        }
        
        .gradio-container .primary {
            background-color: #28a745 !important;
            border-color: #28a745 !important;
        }
        
        .gradio-container .primary:hover {
            background-color: #218838 !important;
        }
        
        .chat-container {
            height: 100px;
            width: 100%;
            max-width: 900px;
            overflow-y: auto;
            border: 2px solid #404040;
            border-radius: 12px;
            padding: 20px;
            background-color: #1a1a1a;
            margin: 20px auto;
            box-shadow: 0 4px 12px rgba(0,0,0,0.5);
        }
        
        .message {
            margin-bottom: 20px;
            display: flex;
            align-items: flex-start;
            gap: 12px;
        }
        
        .user-message {
            flex-direction: row-reverse;
        }
        
        .user-message .message-content {
            background-color: #333333;
            color: #ffffff;
            border-radius: 18px 18px 6px 18px;
            padding: 12px 16px;
            max-width: 70%;
            margin-left: auto;
            box-shadow: 0 2px 8px rgba(0,0,0,0.4);
            border: 1px solid #555;
        }
        
        .ai-message .message-content {
            background-color: #2d2d2d;
            color: #ffffff;
            border-radius: 18px 18px 18px 6px;
            padding: 12px 16px;
            max-width: 70%;
            box-shadow: 0 2px 8px rgba(0,0,0,0.4);
            border: 1px solid #404040;
        }
        
        .avatar {
            width: 32px;
            height: 32px;
            border-radius: 50%;
            background-color: #28a745;
            display: flex;
            align-items: center;
            justify-content: center;
            font-size: 16px;
            flex-shrink: 0;
            box-shadow: 0 2px 6px rgba(40,167,69,0.3);
        }
        
        .message-image {
            max-width: 300px;
            max-height: 400px;
            border-radius: 8px;
            margin-bottom: 8px;
            display: block;
            border: 1px solid #555;
        }
        
        .message-text {
            line-height: 1.4;
            white-space: pre-wrap;
            font-size: 15px;
        }
        
        .user-message .message-text {
            text-align: left;
        }
        
        /* Scrollbar styling for dark theme */
        .chat-container::-webkit-scrollbar {
            width: 8px;
        }
        
        .chat-container::-webkit-scrollbar-track {
            background: #2d2d2d;
            border-radius: 4px;
        }
        
        .chat-container::-webkit-scrollbar-thumb {
            background: #555;
            border-radius: 4px;
        }
        
        .chat-container::-webkit-scrollbar-thumb:hover {
            background: #777;
        }
        
        .file-attachment {
            display: inline-flex;
            align-items: center;
            gap: 8px;
            background-color: #404040;
            border: 1px solid #555;
            border-radius: 8px;
            padding: 8px 12px;
            margin-bottom: 8px;
            font-size: 14px;
        }
        
        .file-icon {
            font-size: 16px;
        }
        
        .file-name {
            color: #ffffff;
            font-weight: 500;
        }
        
        .user-message .file-attachment {
            background-color: #555;
            border-color: #666;
        }
    </style>
    """
    
    with gr.Blocks(
        title="AI Doctor - Consistent Dark Theme",
        theme=gr.themes.Base().set(
            body_background_fill='#1a1a1a',
            body_background_fill_dark='#1a1a1a',
            block_background_fill='#1a1a1a',
            block_background_fill_dark='#1a1a1a',
            border_color_primary='#2d2d2d',
            border_color_primary_dark='#2d2d2d'
        ),
        css=custom_css
    ) as interface:
        
        gr.Markdown(
            """
            # 🩺 AI Medical Consultant
            """,
            # elem_classes="dark-markdown"
        )
        
        # Chat display area
        chat_html = gr.HTML(
            value=f"""
            {custom_css}
            <div class="chat-container" id="chat-container">
                <div class="message ai-message">
                    <div class="avatar">🩺</div>
                    <div class="message-content">
                        <div class="message-text">Hello! I'm your AI medical consultant. You can ask me questions about health concerns and upload medical images for analysis. How can I help you today?</div>
                    </div>
                </div>
            </div>
            """,
            elem_id="chat_display"
        )
        
        with gr.Row():
            with gr.Column(scale=3):
                # Message input
                message_input = gr.Textbox(
                    placeholder="Type your medical question here...",
                    label="Your Message",
                    lines=2,
                    max_lines=4
                )
                
                with gr.Row():
                    image_input = gr.Image(
                        type="filepath",
                        label="📷 Upload Medical Image",
                        height=200
                    )
                    audio_input = gr.Audio(
                        sources=["microphone"],
                        type="filepath",
                        label="🎤 Voice Input",
                    )
                    file_input = gr.File(
                        type="filepath", 
                        label="📄 Upload Report/File",
                        height=200
                    )
                
                with gr.Row():
                    send_btn = gr.Button("Send Message", variant="primary", scale=2)
                    clear_btn = gr.Button("Clear Chat", variant="secondary", scale=1)
            
            with gr.Column(scale=1):
                gr.Markdown("### 🔊 Audio Response")
                audio_output = gr.Audio(
                    label="Doctor's Voice Response",
                    autoplay=True
                )
                
                gr.Markdown(
                    """
                    ### 💡 Features:
                    - 📝 Text conversations
                    - 🖼️ **Images display inline** 
                    - 🎤 Voice input support
                    - 🔊 Audio responses
                    - 💭 Conversation memory
                    - 🌙 **Consistent dark theme**
                    """,
                    elem_classes="dark-markdown"
                )
        
        # Event handlers
        def handle_send(message, current_html, audio, image, file):
            return ai_doctor.process_message(message, current_html, audio, image, file)
        
        def handle_clear():
            result = ai_doctor.clear_conversation()
            return result[0], f"""
            {custom_css}
            <div class="chat-container" id="chat-container">
                <div class="message ai-message">
                    <div class="avatar">🩺</div>
                    <div class="message-content">
                        <div class="message-text">Hello! I'm your AI medical consultant. You can ask me questions about health concerns and upload medical images for analysis. How can I help you today?</div>
                    </div>
                </div>
            </div>
            """
        
        # Bind events
        send_btn.click(
            fn=handle_send,
            inputs=[message_input, chat_html, audio_input, image_input,file_input],
            outputs=[message_input, chat_html, audio_output]
        ).then(
            lambda: [None, None, None],
            outputs=[audio_input, image_input, file_input]
        )
        
        message_input.submit(
            fn=handle_send,
            inputs=[message_input, chat_html, audio_input, image_input, file_input],
            outputs=[message_input, chat_html, audio_output]
        ).then(
            lambda: [None, None, None],
            outputs=[audio_input, image_input, file_input]
        )
        
        clear_btn.click(
            fn=handle_clear,
            outputs=[audio_output, chat_html]
        )
        
        # Auto-scroll JavaScript
        interface.load(
            None,
            None,
            None,
            js="""
            function() {
                setInterval(function() {
                    let container = document.getElementById('chat-container');
                    if (container) {
                        container.scrollTop = container.scrollHeight;
                    }
                }, 500);
            }
            """
        )
    
    return interface


if __name__ == "__main__":
    print("🚀 Launching Dark Theme AI Doctor Interface...")
    
    # Check for required environment variables
    if not os.environ.get("GROQ_API_KEY"):
        print("❌ Missing GROQ_API_KEY environment variable!")
        exit(1)
    else:
        print("✅ GROQ_API_KEY: Set")
    
    interface = create_interface()
    interface.launch(
        debug=True,
        show_error=True,
        # server_name="0.0.0.0",
        # server_port=7860,
        # share=True
    )
